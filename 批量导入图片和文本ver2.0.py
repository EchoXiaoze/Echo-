from pptx import Presentation
import pandas as pd
import xlrd
import openpyxl
数据 = pd.read_excel('C:/py/档案.xlsx')
list = []
num_行号=0
mading = -1
for 行号 in 数据['序号']:#统计一共多少学生
    num_行号 = num_行号+1
print('一共',num_行号,'个学生')
for 行号_2 in 数据['序号']:
    
    姓名 = 数据.iloc[行号_2-1,1]
    if 姓名 == '结束':
        break
    else:
        pass
    root = f'C:/py/img/{姓名}/'
    print('正在制作',姓名,'的PPT')
    mading = mading + 1
    n_2 = 0
    n_2 = n_2-1
    list = []
    for 列数 in range(1,数据.shape[1]):
        try:
            list.append(数据.iloc[mading,列数])
        except IndexError:
            print('已遇到‘结束’')
        else:
            pass
 
    #对整篇PPT遍历
    文件 = Presentation('C:/py/beta1.0.pptx')
    total_page = len(文件.slides)
    for i in range(total_page):
        list_pic = []
        list_text = []
        list_text_loc = []
        print('''#####'''
            '第',i+1,'页'
            '''#####''')
        页 = 文件.slides[i]
        if str(bool(页.placeholders)) == 'False':#判断本页有无占位符
            print('没有占位符')
            continue
        elif str(bool(页.placeholders)) == 'True':
    # 图片导入
            print('''
            ##
            开始图片导入
            ##
            ''')    
            for 占位符 in 页.placeholders:
                信息 = 占位符.placeholder_format
                占位符类型 = 占位符.name[:5]
                if 占位符类型 == '图片占位符':
                    list_pic.append(f'{信息.idx}')
                    print('发现一个【',占位符类型,'】')
                elif 占位符类型 == '文本占位符':
                    continue
        #循环图片添加
            n=0
            for 索引 in list_pic:
                n=n+1
                for z in range(0,n):
                    页数 = str(i+1)
                    数量 = str(z+1)
                    图片 = 页数 + '-' + 数量 + '.jpg'
                    图片路径 = str(root)+str(图片)
                插入图片位置 = 页.placeholders[int(索引)]
                print('正在导入【',图片,'】')
                print('图片路径',f'{图片路径}')
                try:
                    插入图片 = 插入图片位置.insert_picture(f'{图片路径}')
                except AttributeError:
                    print('老错误')
                print('导入完毕')
            print('''
            ##
            本页没有图片可以导入了
            ##
            ''')
    #文本导入
            print('''
            ##
            开始文本导入
            ##
            ''')
            
            if str(bool(页.placeholders)) == 'False':
                print('没有占位符')
                continue
            elif str(bool(页.placeholders)) == 'True':
                a1 = 0
                for 占位符 in 页.placeholders:#对接页面遍历
                    信息 = 占位符.placeholder_format
                    占位符类型 = 占位符.name[:5]
                    占位 = 0
                    if 占位符类型 == '文本占位符':
                        a1 = a1 + 1
                        list_text.append(f'{信息.idx}')#单个索引的添加
                        list_text_loc.append(页.placeholders[信息.idx])#占位符的位置
                for 导入文本 in range(0,len(list_text_loc)):
                    n_2 = n_2+1
                    try:
                        list_text_loc[导入文本].text = list[n_2]
                    except IndexError:
                        print('已遇到‘结束’')
                    else:
                        list_text_loc[导入文本].text = list[n_2]
    文件.save(f'{姓名}.pptx')